import os
import sys
import sip
from PyQt4 import QtGui, QtCore, uic
from PyQt4.QtGui import QApplication, QMainWindow, QMessageBox
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
#from logic import build_deck

class MainWindow(QMainWindow):

    INSTRUCTION_TEXT = 'Copy the entire HW folder (and make sure the folder is called HW) to the AUTO_HW folder before pressing the "Build Slides" button!'

    def __init__(self):
        super(MainWindow, self).__init__()
        self.ui = uic.loadUi('imgs/UI.ui')
        self.ui.setWindowIcon(QtGui.QIcon("imgs/pencil.ico"))

        self.ui.progress_bar.setValue(0)

        self.ui.text_output.setPlainText(self.INSTRUCTION_TEXT)

        self.ui.start_btn.clicked.connect(self.build_deck)

        self.ui.show()

    def get_list_of_names(self, start_dir, country_list):
        names_from_files = []

        for x in country_list:
            country_path = start_dir + x
            file_list = [x for x in os.listdir(country_path)]
            for i in file_list:
                line = i.split("-")
                if i[-3:] == "jpg":
                    line2 = line[0].split("_")
                    if line2[1] not in names_from_files:
                        names_from_files.append(line2[1])

        return names_from_files

    def build_deck(self) :

        start_dir = os.getcwd() + "\HW\\"
        country_list = [x for x in os.listdir(start_dir)]
        names_list = self.get_list_of_names(start_dir, country_list)
        HW = {}

        total = len(country_list) * len(names_list)
        self.ui.progress_bar.setMaximum(total)


        for c in country_list:
            HW[c] = {}
            for n in names_list:
                HW[c][n] = []

        # Creates a dictionary mapping out countries, names, and file names for each

        for c in HW:
            for n in HW[c]:
                count = 1
                c_range = 5
                if c == "USA" or c == "EU":
                    c_range = 10
                for x in range(0, c_range):
                    if count == 10:
                        file_name = c + "_" + n + "-10.jpg"
                    else:
                        file_name = c + "_" + n + "-0" + str(count) + ".jpg"
                    HW[c][n].append(file_name)
                    count += 1



        prs_slides = {}

        prs = Presentation('ppts/aw_template.pptx') # Opens AW master template
        hw_slide_layout = prs.slide_layouts[5] # Slide used for handwriting samples

        ERROR_IMG = "imgs/error_img.png"
        bar_val = 0
        out_text = ""
        for name in names_list:
            for country in country_list:

                count = 0

                out_text = "{} {} slide added\n".format(name, country)
                self.ui.text_output.setPlainText(out_text)
                self.ui.progress_bar.setValue(bar_val)
                sld = "Slide %d" % (x)
                prs_slides[sld] = prs.slides.add_slide(hw_slide_layout)
                title = prs_slides[sld].placeholders[18] #Title text in Red
                subtitle = prs_slides[sld].placeholders[16] # Subtitle text in Grey

                title.text = "HANDWRITING SAMPLES"
                subtitle.text = name
                textbox = prs_slides[sld].shapes[2]
                sp = textbox.element
                sp.getparent().remove(sp)

                ### Add grey box

                left = Inches(1.25)
                top = Inches(1.5)
                ht = Inches(4.5)
                wd = Inches(7.5)
                prs_slides[sld].shapes.add_picture("imgs/grey_box.png", left, top, wd, ht)

                ### Add flag image

                flag_img = "imgs/flag_{}.png".format(country)
                left = Inches(8)
                top = Inches(0.7)
                ht = Inches(0.7)
                wd = Inches(0.7)
                try:
                    prs_slides[sld].shapes.add_picture(flag_img, left, top, wd, ht)
                except:
                    prs_slides[sld].shapes.add_picture(ERROR_IMG, left, top, wd, ht)
                ### Add HW samples to each
                for x in range(len(HW[country][name])):
                    img_file = "HW\\" + country + "\\" + HW[country][name][x]
                    try:
                        im = Image.open(img_file)
                    except:
                        im = Image.open(ERROR_IMG)
                    ht = Inches(1.25)
                    wd = Inches(2.0)
                    #print img_file
                    # im.show()
                    # raw_input("Enter")

                    if "01" in img_file or "06" in img_file or "02" in img_file:
                        top = Inches(1.75)
                    if "07" in img_file or "03" in img_file or "08" in img_file:
                        top = Inches(3.125)
                    if "04" in img_file or "09" in img_file or "05" in img_file or "10" in img_file:
                        top = Inches(4.5)

                    if "01" in img_file or "07" in img_file or "04" in img_file:
                        left = Inches(1.5)
                    if "06" in img_file or "03" in img_file or "09" in img_file:
                        left = Inches(4.0)
                    if "02" in img_file or "08" in img_file or "05" in img_file or "10" in img_file:
                        left = Inches(6.5)
                    try:
                        prs_slides[sld].shapes.add_picture(img_file, left, top, wd, ht)
                    except:
                        prs_slides[sld].shapes.add_picture(ERROR_IMG, left, top, wd, ht)

                bar_val += 1
                QApplication.processEvents()
        try:
            prs.save('hw_slides.pptx')
            self.open_deck('hw_slides.pptx')
        except IOError:
            print "File already open"
            QMessageBox.warning(self, "File already open", "The hw_slides file is currently open; close this file first and try again")
        print "DONE"
        self.ui.progress_bar.setValue(total)
        self.ui.text_output.setPlainText("All slides have been built!")


    def open_deck(self, file_name):
        try:
            os.system("start " + str(file_name))
        except:
            print file_name, "not found"


if __name__ == '__main__':
    app = QApplication(sys.argv)
    app.setStyle("Plastique")
    app.setWindowIcon(QtGui.QIcon("book.ico"))
    window = MainWindow()

    sip.setdestroyonexit(False)
    sys.exit(app.exec_())